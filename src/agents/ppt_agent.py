"""
McLeuker Agentic AI Platform - PowerPoint Execution Agent

Generates PowerPoint presentations from structured data.
"""

import os
from typing import Optional, List, Dict, Any
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from src.agents.base_agent import BaseExecutionAgent
from src.models.schemas import (
    StructuredOutput,
    StructuredPresentation,
    GeneratedFile,
    OutputFormat
)


class PPTAgent(BaseExecutionAgent):
    """
    PowerPoint Execution Agent
    
    Generates PowerPoint presentations from structured presentation data.
    """
    
    # Color scheme
    COLORS = {
        'primary': RgbColor(31, 78, 121),      # Dark blue
        'secondary': RgbColor(46, 117, 182),   # Medium blue
        'accent': RgbColor(68, 114, 196),      # Light blue
        'text_dark': RgbColor(51, 51, 51),     # Dark gray
        'text_light': RgbColor(255, 255, 255), # White
        'background': RgbColor(245, 245, 245)  # Light gray
    }
    
    @property
    def supported_formats(self) -> List[OutputFormat]:
        return [OutputFormat.PPT]
    
    async def execute(
        self,
        structured_output: StructuredOutput,
        filename_prefix: Optional[str] = None
    ) -> List[GeneratedFile]:
        """
        Generate PowerPoint files from structured output.
        
        Args:
            structured_output: The structured data containing presentations
            filename_prefix: Optional prefix for filenames
            
        Returns:
            List[GeneratedFile]: List of generated files
        """
        generated_files = []
        
        if not structured_output.presentations:
            return generated_files
        
        for presentation in structured_output.presentations:
            ppt_file = await self._generate_ppt(presentation, filename_prefix)
            if ppt_file:
                generated_files.append(ppt_file)
        
        return generated_files
    
    async def _generate_ppt(
        self,
        presentation: StructuredPresentation,
        filename_prefix: Optional[str] = None
    ) -> Optional[GeneratedFile]:
        """Generate a PowerPoint file from a structured presentation."""
        try:
            filename = self._generate_filename(filename_prefix, "pptx")
            filepath = self._get_file_path(filename)
            
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
            prs.slide_height = Inches(7.5)
            
            # Add title slide
            self._add_title_slide(prs, presentation.title)
            
            # Add content slides
            for slide_data in presentation.slides:
                layout = slide_data.get('layout', 'content')
                
                if layout == 'title':
                    self._add_section_slide(prs, slide_data)
                elif layout == 'two_column':
                    self._add_two_column_slide(prs, slide_data)
                elif layout == 'chart':
                    self._add_chart_slide(prs, slide_data)
                elif layout == 'image':
                    self._add_image_slide(prs, slide_data)
                elif layout == 'blank':
                    self._add_blank_slide(prs, slide_data)
                else:
                    self._add_content_slide(prs, slide_data)
            
            # Add thank you slide
            self._add_thank_you_slide(prs)
            
            # Save presentation
            prs.save(filepath)
            
            return self._create_generated_file(filename, OutputFormat.PPT)
            
        except Exception as e:
            print(f"Error generating PowerPoint file: {e}")
            return None
    
    def _add_title_slide(self, prs: Presentation, title: str):
        """Add a title slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background shape
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['primary']
        background.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(2)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['text_light']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add subtitle with date
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5),
            Inches(12.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = f"Generated by McLeuker AI | {datetime.now().strftime('%B %Y')}"
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = self.COLORS['text_light']
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    def _add_content_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Add a content slide with title and bullet points."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS['primary']
        title_bar.line.fill.background()
        
        # Add title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get('title', 'Slide Title')
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['text_light']
        
        # Add content
        content_top = Inches(1.5)
        
        # Add bullet points if present
        bullet_points = slide_data.get('bullet_points', [])
        if bullet_points:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), content_top,
                Inches(11.833), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for i, point in enumerate(bullet_points):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()
                
                para.text = f"• {point}"
                para.font.size = Pt(20)
                para.font.color.rgb = self.COLORS['text_dark']
                para.space_after = Pt(12)
        
        # Add content text if present
        content = slide_data.get('content', '')
        if content and not bullet_points:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), content_top,
                Inches(11.833), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            para = content_frame.paragraphs[0]
            para.text = content
            para.font.size = Pt(18)
            para.font.color.rgb = self.COLORS['text_dark']
    
    def _add_section_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Add a section divider slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['secondary']
        background.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3),
            Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get('title', 'Section')
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['text_light']
        title_para.alignment = PP_ALIGN.CENTER
    
    def _add_two_column_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Add a two-column slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS['primary']
        title_bar.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get('title', 'Comparison')
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['text_light']
        
        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(5.9), Inches(5.5)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        left_para = left_frame.paragraphs[0]
        left_para.text = slide_data.get('left_content', 'Left column content')
        left_para.font.size = Pt(16)
        
        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(6.9), Inches(1.5),
            Inches(5.9), Inches(5.5)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        right_para = right_frame.paragraphs[0]
        right_para.text = slide_data.get('right_content', 'Right column content')
        right_para.font.size = Pt(16)
    
    def _add_chart_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Add a chart placeholder slide."""
        # For now, add a content slide with a note about charts
        self._add_content_slide(prs, {
            'title': slide_data.get('title', 'Chart'),
            'content': slide_data.get('content', 'Chart visualization would be displayed here.')
        })
    
    def _add_image_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Add an image placeholder slide."""
        self._add_content_slide(prs, {
            'title': slide_data.get('title', 'Visual'),
            'content': slide_data.get('content', 'Image would be displayed here.')
        })
    
    def _add_blank_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Add a blank slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
    
    def _add_thank_you_slide(self, prs: Presentation):
        """Add a thank you/closing slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['primary']
        background.line.fill.background()
        
        # Add thank you text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(2)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Thank You"
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['text_light']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add contact info
        contact_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5),
            Inches(12.333), Inches(1)
        )
        contact_frame = contact_box.text_frame
        contact_para = contact_frame.paragraphs[0]
        contact_para.text = "Powered by McLeuker AI"
        contact_para.font.size = Pt(18)
        contact_para.font.color.rgb = self.COLORS['text_light']
        contact_para.alignment = PP_ALIGN.CENTER
    
    async def generate_simple_presentation(
        self,
        title: str,
        slides: List[Dict[str, Any]],
        filename_prefix: Optional[str] = None
    ) -> Optional[GeneratedFile]:
        """
        Generate a simple presentation.
        
        Args:
            title: Presentation title
            slides: List of slide data dictionaries
            filename_prefix: Optional prefix for filename
            
        Returns:
            GeneratedFile or None if failed
        """
        presentation = StructuredPresentation(
            title=title,
            slides=slides,
            theme="professional"
        )
        
        structured_output = StructuredOutput(presentations=[presentation])
        files = await self.execute(structured_output, filename_prefix)
        return files[0] if files else None
