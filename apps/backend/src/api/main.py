"""
McLeuker AI V10.0 - Agentic Execution Platform
=====================================================
True end-to-end agentic AI with multi-layer reasoning,
execution engine, domain agents, and web automation.

Route Structure:
  /api/v1/chat/stream      — SSE streaming chat (research mode)
  /api/v1/generate-file    — File generation (Excel, PDF, etc.)
  /api/v2/execute/stream   — SSE streaming execution (agent mode)
  /api/chat/stream         — Legacy alias
  /api/execute/stream      — Legacy alias
"""

import os
import json
import uuid
import base64
import asyncio
import logging
import re
from datetime import datetime
from typing import Optional, AsyncGenerator, List, Union
from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse, StreamingResponse
from pydantic import BaseModel, Field

from src.core.reasoning_orchestrator import ReasoningOrchestrator
from src.services.nano_banana import nano_banana_service

logger = logging.getLogger(__name__)

# Enhancement imports
try:
    from src.enhancement.execution import ExecutionEngine
    from src.enhancement import McLeukerEnhancement, EnhancementConfig
    ENHANCEMENT_AVAILABLE = True
except ImportError as e:
    logger.warning(f"Enhancement modules not available: {e}")
    ENHANCEMENT_AVAILABLE = False

# Initialize orchestrator
reasoning_orchestrator = ReasoningOrchestrator()

# Initialize execution engine
execution_engine = None
enhancement = None
if ENHANCEMENT_AVAILABLE:
    try:
        execution_engine = ExecutionEngine()
        enhancement = McLeukerEnhancement(
            llm_client=None,
            config=EnhancementConfig(
                enable_domain_agents=True,
                enable_tool_stability=True,
                enable_file_analysis=True,
                enable_execution_engine=True,
            )
        )
    except Exception as e:
        logger.warning(f"Enhancement initialization failed: {e}")


# =============================================================================
# FastAPI App
# =============================================================================

app = FastAPI(
    title="McLeuker AI V10.0",
    description="Agentic Execution Platform — Multi-Layer Reasoning, Web Automation, Domain Agents",
    version="10.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# =============================================================================
# Request / Response Models
# =============================================================================

class ChatRequest(BaseModel):
    """Flexible chat request — accepts both single message and messages array."""
    message: Optional[str] = None
    messages: Optional[list] = None
    session_id: Optional[str] = None
    mode: str = "quick"
    sector: Optional[str] = None
    stream: Optional[bool] = True
    enable_tools: Optional[bool] = True
    user_id: Optional[str] = None
    conversation_id: Optional[str] = None
    attachments: Optional[List[str]] = None

    def get_query(self) -> str:
        """Extract the user query from either message or messages array."""
        if self.message:
            return self.message
        if self.messages:
            for msg in reversed(self.messages):
                if isinstance(msg, dict) and msg.get("role") == "user":
                    content = msg.get("content", "")
                    if isinstance(content, str):
                        return content
                    if isinstance(content, list):
                        # Handle multimodal content array
                        texts = [p.get("text", "") for p in content if isinstance(p, dict) and p.get("type") == "text"]
                        return " ".join(texts)
            # Fallback: last message content
            last = self.messages[-1]
            if isinstance(last, dict):
                return str(last.get("content", ""))
            return str(last)
        return ""


class ExecuteRequest(BaseModel):
    """Request for agent execution mode."""
    task: Optional[str] = None
    query: Optional[str] = None
    session_id: Optional[str] = None
    user_id: Optional[str] = None
    conversation_id: Optional[str] = None
    mode: str = "agent"
    context: Optional[dict] = None

    def get_query(self) -> str:
        return self.task or self.query or ""


class GenerateFileRequest(BaseModel):
    """Request to generate a file (Excel, PDF, Word, etc.)."""
    content: Optional[str] = None
    prompt: Optional[str] = None
    title: str = "McLeuker AI Report"
    file_type: str = "excel"  # excel, pdf, word, markdown, csv, pptx
    data: Optional[list] = None  # Structured data for Excel


class ImageGenerateRequest(BaseModel):
    prompt: str
    style: Optional[str] = None
    aspect_ratio: str = "1:1"


class ImageEditRequest(BaseModel):
    image_data: str
    edit_prompt: str
    mime_type: str = "image/png"


class ImageAnalyzeRequest(BaseModel):
    image_data: str
    question: str
    mime_type: str = "image/png"


class DocumentGenerateRequest(BaseModel):
    content: str
    title: str
    format: str
    template: Optional[str] = None


class ExportChatRequest(BaseModel):
    messages: List[dict]
    title: str
    format: str


# =============================================================================
# Health & Root
# =============================================================================

@app.get("/health")
async def health_check():
    return {
        "status": "healthy",
        "version": "10.0.0",
        "timestamp": datetime.now().isoformat(),
        "features": {
            "multi_layer_reasoning": True,
            "streaming": True,
            "realtime_data": True,
            "dynamic_sources": True,
            "nano_banana": True,
            "file_upload": True,
            "file_generation": True,
            "execution_engine": execution_engine is not None,
            "domain_agents": enhancement is not None,
            "web_automation": True,
            "file_analysis": True,
        }
    }


@app.get("/")
async def root():
    return {
        "name": "McLeuker AI",
        "version": "10.0.0",
        "description": "Agentic Execution Platform — Multi-Layer Reasoning, Web Automation, Domain Agents",
        "endpoints": {
            "health": "/health",
            "chat_stream_v1": "/api/v1/chat/stream",
            "execute_stream_v2": "/api/v2/execute/stream",
            "generate_file": "/api/v1/generate-file",
            "image_generate": "/api/image/generate",
            "image_edit": "/api/image/edit",
            "image_analyze": "/api/image/analyze",
            "file_upload": "/api/upload",
            "file_download": "/api/files/{file_id}",
            "document_generate": "/api/document/generate",
            "export_chat": "/api/export/chat",
            "capabilities": "/api/capabilities",
        }
    }


# =============================================================================
# File Generation Utility
# =============================================================================

def _generate_excel_from_content(content: str, title: str) -> tuple:
    """Generate an Excel file from content. Returns (file_path, filename)."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    os.makedirs("/tmp/mcleuker_files", exist_ok=True)
    file_id = str(uuid.uuid4())[:12]

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]

    # Style definitions
    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="6B21A8", end_color="6B21A8", fill_type="solid")
    header_alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    data_alignment = Alignment(vertical="center", wrap_text=True)
    thin_border = Border(
        left=Side(style="thin", color="CCCCCC"),
        right=Side(style="thin", color="CCCCCC"),
        top=Side(style="thin", color="CCCCCC"),
        bottom=Side(style="thin", color="CCCCCC"),
    )
    alt_fill = PatternFill(start_color="F5F0FF", end_color="F5F0FF", fill_type="solid")

    # Parse content into rows
    lines = content.strip().split('\n')
    rows = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        # Skip markdown table separators
        if re.match(r'^[\|\-\+\s:]+$', line):
            continue
        # Handle pipe-delimited (markdown table)
        if '|' in line:
            cells = [c.strip() for c in line.split('|') if c.strip()]
        elif '\t' in line:
            cells = [c.strip() for c in line.split('\t')]
        elif ',' in line and line.count(',') >= 2:
            cells = [c.strip().strip('"') for c in line.split(',')]
        else:
            cells = [line]
        if cells:
            rows.append(cells)

    # Write data
    for row_idx, cells in enumerate(rows, 1):
        for col_idx, cell_value in enumerate(cells, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
            cell.border = thin_border
            if row_idx == 1:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_alignment
            else:
                cell.alignment = data_alignment
                if row_idx % 2 == 0:
                    cell.fill = alt_fill

    # Auto-adjust column widths
    for column in ws.columns:
        max_length = 0
        column_letter = column[0].column_letter
        for cell in column:
            try:
                if cell.value and len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        ws.column_dimensions[column_letter].width = min(max(max_length + 4, 12), 50)

    # Freeze header row
    ws.freeze_panes = "A2"

    file_path = f"/tmp/mcleuker_files/{file_id}.xlsx"
    wb.save(file_path)
    filename = f"{title}.xlsx"
    return file_path, filename, file_id


def _generate_csv_from_content(content: str, title: str) -> tuple:
    """Generate a CSV file from content."""
    import csv
    import io

    os.makedirs("/tmp/mcleuker_files", exist_ok=True)
    file_id = str(uuid.uuid4())[:12]

    lines = content.strip().split('\n')
    rows = []
    for line in lines:
        line = line.strip()
        if not line or re.match(r'^[\|\-\+\s:]+$', line):
            continue
        if '|' in line:
            cells = [c.strip() for c in line.split('|') if c.strip()]
        elif '\t' in line:
            cells = [c.strip() for c in line.split('\t')]
        else:
            cells = [line]
        if cells:
            rows.append(cells)

    file_path = f"/tmp/mcleuker_files/{file_id}.csv"
    with open(file_path, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        for row in rows:
            writer.writerow(row)

    filename = f"{title}.csv"
    return file_path, filename, file_id


def _generate_pdf_from_content(content: str, title: str) -> tuple:
    """Generate a PDF file from content."""
    from fpdf import FPDF

    os.makedirs("/tmp/mcleuker_files", exist_ok=True)
    file_id = str(uuid.uuid4())[:12]

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", "B", 16)
    pdf.cell(0, 10, title, ln=True, align="C")
    pdf.ln(10)
    pdf.set_font("Arial", "", 11)

    for line in content.split('\n'):
        line = line.strip()
        if line.startswith('# '):
            pdf.set_font("Arial", "B", 14)
            pdf.multi_cell(0, 8, line[2:])
            pdf.set_font("Arial", "", 11)
        elif line.startswith('## '):
            pdf.set_font("Arial", "B", 12)
            pdf.multi_cell(0, 8, line[3:])
            pdf.set_font("Arial", "", 11)
        elif line.startswith('- ') or line.startswith('* '):
            pdf.multi_cell(0, 7, f"  \u2022 {line[2:]}")
        else:
            pdf.multi_cell(0, 7, line)

    file_path = f"/tmp/mcleuker_files/{file_id}.pdf"
    pdf.output(file_path)
    filename = f"{title}.pdf"
    return file_path, filename, file_id


def _generate_word_from_content(content: str, title: str) -> tuple:
    """Generate a Word document from content."""
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    os.makedirs("/tmp/mcleuker_files", exist_ok=True)
    file_id = str(uuid.uuid4())[:12]

    doc = Document()
    heading = doc.add_heading(title, 0)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

    paragraphs = content.split('\n\n')
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if para.startswith('# '):
            doc.add_heading(para[2:], level=1)
        elif para.startswith('## '):
            doc.add_heading(para[3:], level=2)
        elif para.startswith('### '):
            doc.add_heading(para[4:], level=3)
        elif para.startswith('- ') or para.startswith('* '):
            for item in para.split('\n'):
                item = item.strip()
                if item.startswith('- ') or item.startswith('* '):
                    doc.add_paragraph(item[2:], style='List Bullet')
        else:
            doc.add_paragraph(para)

    file_path = f"/tmp/mcleuker_files/{file_id}.docx"
    doc.save(file_path)
    filename = f"{title}.docx"
    return file_path, filename, file_id


def _generate_markdown_from_content(content: str, title: str) -> tuple:
    """Generate a Markdown file from content."""
    os.makedirs("/tmp/mcleuker_files", exist_ok=True)
    file_id = str(uuid.uuid4())[:12]

    full_content = f"# {title}\n\n{content}"
    file_path = f"/tmp/mcleuker_files/{file_id}.md"
    with open(file_path, 'w', encoding='utf-8') as f:
        f.write(full_content)

    filename = f"{title}.md"
    return file_path, filename, file_id


def _generate_pptx_from_content(content: str, title: str) -> tuple:
    """Generate a PowerPoint file from content."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    os.makedirs("/tmp/mcleuker_files", exist_ok=True)
    file_id = str(uuid.uuid4())[:12]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_content = content.split('---SLIDE---') if '---SLIDE---' in content else [content]

    for slide_content in slides_content:
        slide_content = slide_content.strip()
        if not slide_content:
            continue
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        lines = slide_content.split('\n')
        title_text = lines[0] if lines else title
        body_text = '\n'.join(lines[1:]) if len(lines) > 1 else ""
        slide.shapes.title.text = title_text
        if len(slide.shapes.placeholders) > 1:
            slide.shapes.placeholders[1].text = body_text

    file_path = f"/tmp/mcleuker_files/{file_id}.pptx"
    prs.save(file_path)
    filename = f"{title}.pptx"
    return file_path, filename, file_id


# =============================================================================
# V1 Chat Streaming Endpoint (Research Mode)
# =============================================================================

async def _chat_stream_generator(request: ChatRequest) -> AsyncGenerator[str, None]:
    """
    SSE generator for chat streaming with multi-layer reasoning.
    Also detects file generation requests and emits download events.
    """
    session_id = request.session_id or request.conversation_id or str(uuid.uuid4())
    query = request.get_query()

    if not query:
        yield f"data: {json.dumps({'type': 'error', 'data': {'message': 'No message provided'}})}\n\n"
        return

    try:
        # Yield start event
        yield f"data: {json.dumps({'type': 'start', 'data': {'session_id': session_id}})}\n\n"

        full_content = ""
        all_sources = []
        all_layers = []
        credits_used = 2

        async for event in reasoning_orchestrator.process_with_reasoning(
            query=query,
            session_id=session_id,
            mode=request.mode
        ):
            event_type = event.get("type")
            event_data = event.get("data", {})

            # Forward all events to frontend
            yield f"data: {json.dumps(event)}\n\n"

            # Track content for file generation detection
            if event_type == "content":
                full_content += event_data.get("chunk", "")
            elif event_type == "complete":
                final = event_data.get("content", "")
                if final and len(final) > len(full_content):
                    full_content = final
                credits_used = event_data.get("credits_used", 2)
                all_sources = event_data.get("sources", all_sources)

            if event_type in ["layer_start", "layer_complete"]:
                await asyncio.sleep(0.1)
            elif event_type == "sub_step":
                await asyncio.sleep(0.05)

        # ── Post-processing: Auto-detect if content contains tabular data ──
        # If the response contains table-like data, auto-generate an Excel file
        table_lines = [l for l in full_content.split('\n') if '|' in l and not l.strip().startswith('---')]
        if len(table_lines) >= 3:
            try:
                # Extract the table portion
                table_content = '\n'.join(table_lines)
                file_path, filename, file_id = _generate_excel_from_content(table_content, "McLeuker AI Data")

                yield f"data: {json.dumps({'type': 'download', 'data': {'filename': filename, 'download_url': f'/api/files/{file_id}.xlsx', 'file_id': f'{file_id}.xlsx', 'file_type': 'excel'}})}\n\n"
            except Exception as e:
                logger.warning(f"Auto Excel generation failed: {e}")

    except Exception as e:
        logger.error(f"Chat stream error: {e}", exc_info=True)
        yield f"data: {json.dumps({'type': 'error', 'data': {'message': str(e)}})}\n\n"


@app.post("/api/v1/chat/stream")
async def chat_stream_v1(request: ChatRequest):
    """V1 SSE streaming chat endpoint — used by research mode."""
    return StreamingResponse(
        _chat_stream_generator(request),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        }
    )


# Legacy alias
@app.post("/api/chat/stream")
async def chat_stream_legacy(request: ChatRequest):
    """Legacy alias for /api/v1/chat/stream."""
    return StreamingResponse(
        _chat_stream_generator(request),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        }
    )


@app.post("/api/chat")
async def chat(request: ChatRequest):
    """Non-streaming chat endpoint."""
    session_id = request.session_id or str(uuid.uuid4())
    query = request.get_query()

    try:
        reasoning_layers = []
        final_content = ""
        sources = []
        credits_used = 2

        async for event in reasoning_orchestrator.process_with_reasoning(
            query=query,
            session_id=session_id,
            mode=request.mode
        ):
            event_type = event.get("type")
            event_data = event.get("data", {})

            if event_type == "layer_start":
                reasoning_layers.append({
                    "id": event_data.get("layer_id"),
                    "layer_num": event_data.get("layer_num"),
                    "type": event_data.get("type"),
                    "title": event_data.get("title"),
                    "sub_steps": [],
                    "status": "active"
                })
            elif event_type == "sub_step":
                layer_id = event_data.get("layer_id")
                for layer in reasoning_layers:
                    if layer["id"] == layer_id:
                        layer["sub_steps"].append({
                            "step": event_data.get("step"),
                            "status": event_data.get("status")
                        })
            elif event_type == "layer_complete":
                layer_id = event_data.get("layer_id")
                for layer in reasoning_layers:
                    if layer["id"] == layer_id:
                        layer["status"] = "complete"
                        layer["content"] = event_data.get("content", "")
            elif event_type == "source":
                sources.append(event_data)
            elif event_type == "content":
                final_content += event_data.get("chunk", "")
            elif event_type == "complete":
                final_content = event_data.get("content", final_content)
                credits_used = event_data.get("credits_used", 2)
            elif event_type == "error":
                return JSONResponse(
                    status_code=500,
                    content={"success": False, "error": event_data.get("message", "Unknown error")}
                )

        return {
            "success": True,
            "response": {"content": final_content, "sources": sources},
            "reasoning_layers": reasoning_layers,
            "credits_used": credits_used,
            "session_id": session_id
        }

    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"success": False, "error": str(e)}
        )


# =============================================================================
# V1 File Generation Endpoint
# =============================================================================

@app.post("/api/v1/generate-file")
async def generate_file_v1(request: GenerateFileRequest):
    """
    Generate a downloadable file from content.
    
    Supports: excel, csv, pdf, word, markdown, pptx
    Returns a download URL that the frontend can use.
    """
    try:
        content = request.content or request.prompt or ""
        title = request.title or "McLeuker AI Report"
        file_type = request.file_type.lower()

        if not content:
            return JSONResponse(
                status_code=400,
                content={"success": False, "error": "No content provided for file generation"}
            )

        if file_type in ("excel", "xlsx"):
            file_path, filename, file_id = _generate_excel_from_content(content, title)
            ext = "xlsx"
        elif file_type == "csv":
            file_path, filename, file_id = _generate_csv_from_content(content, title)
            ext = "csv"
        elif file_type == "pdf":
            file_path, filename, file_id = _generate_pdf_from_content(content, title)
            ext = "pdf"
        elif file_type in ("word", "docx"):
            file_path, filename, file_id = _generate_word_from_content(content, title)
            ext = "docx"
        elif file_type == "markdown":
            file_path, filename, file_id = _generate_markdown_from_content(content, title)
            ext = "md"
        elif file_type in ("pptx", "powerpoint"):
            file_path, filename, file_id = _generate_pptx_from_content(content, title)
            ext = "pptx"
        else:
            return JSONResponse(
                status_code=400,
                content={"success": False, "error": f"Unsupported file type: {file_type}"}
            )

        return {
            "success": True,
            "file_id": f"{file_id}.{ext}",
            "filename": filename,
            "download_url": f"/api/files/{file_id}.{ext}",
            "format": file_type,
        }

    except Exception as e:
        logger.error(f"File generation error: {e}", exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"success": False, "error": str(e)}
        )


# =============================================================================
# V2 Agent Execution Streaming Endpoint
# =============================================================================

async def _execute_stream_generator(request: ExecuteRequest) -> AsyncGenerator[str, None]:
    """
    SSE generator for agent execution mode.
    Emits: execution_start, step_update, execution_progress, execution_reasoning,
           browser_screenshot, execution_artifact, download, execution_complete, execution_error
    """
    query = request.get_query()
    execution_id = str(uuid.uuid4())

    if not query:
        yield f"data: {json.dumps({'type': 'error', 'data': {'error': 'No task provided'}})}\n\n"
        return

    # Emit execution start
    yield f"data: {json.dumps({'type': 'execution_start', 'data': {'execution_id': execution_id, 'task': query}})}\n\n"

    try:
        if execution_engine:
            # Use the V9 execution engine
            step_count = 0
            async for event in execution_engine.execute(
                query=query,
                user_id=request.user_id,
                session_id=request.session_id or request.conversation_id,
                context=request.context,
            ):
                event_type = event.get("type", "")

                # Map execution engine events to frontend-expected events
                if event_type == "planning_start":
                    yield f"data: {json.dumps({'type': 'step_update', 'data': {'step_id': 'planning', 'phase': 'planning', 'title': 'Analyzing and planning task...', 'status': 'active'}})}\n\n"
                elif event_type == "planning_complete":
                    steps = event.get("data", {}).get("steps", [])
                    yield f"data: {json.dumps({'type': 'step_update', 'data': {'step_id': 'planning', 'phase': 'planning', 'title': f'Plan ready — {len(steps)} steps', 'status': 'complete'}})}\n\n"
                    yield f"data: {json.dumps({'type': 'execution_progress', 'data': {'progress': 10, 'status': 'executing'}})}\n\n"
                elif event_type == "step_start":
                    step_count += 1
                    step_data = event.get("data", {})
                    yield f"data: {json.dumps({'type': 'step_update', 'data': {'step_id': f'step-{step_count}', 'phase': 'execution', 'title': step_data.get('description', f'Step {step_count}'), 'status': 'active', 'detail': step_data.get('tool', '')}})}\n\n"
                elif event_type == "step_complete":
                    step_data = event.get("data", {})
                    progress = min(10 + (step_count * 80 // max(step_data.get("total_steps", 5), 1)), 90)
                    yield f"data: {json.dumps({'type': 'step_update', 'data': {'step_id': f'step-{step_count}', 'phase': 'execution', 'title': step_data.get('description', f'Step {step_count} complete'), 'status': 'complete'}})}\n\n"
                    yield f"data: {json.dumps({'type': 'execution_progress', 'data': {'progress': progress}})}\n\n"
                elif event_type == "screenshot":
                    ss_data = event.get("data", {})
                    yield f"data: {json.dumps({'type': 'browser_screenshot', 'data': {'screenshot': ss_data.get('image', ''), 'url': ss_data.get('url', ''), 'title': ss_data.get('title', ''), 'step': step_count, 'action': ss_data.get('action', 'Browsing')}})}\n\n"
                elif event_type == "result":
                    result_data = event.get("data", {})
                    content = result_data.get("content", "")
                    # Emit content
                    yield f"data: {json.dumps({'type': 'content', 'data': {'chunk': content}})}\n\n"
                    # Check for downloads
                    downloads = result_data.get("downloads", [])
                    for dl in downloads:
                        yield f"data: {json.dumps({'type': 'download', 'data': dl})}\n\n"
                elif event_type == "error":
                    yield f"data: {json.dumps({'type': 'execution_error', 'data': event.get('data', {})})}\n\n"
                else:
                    # Forward unknown events
                    yield f"data: {json.dumps(event)}\n\n"

                await asyncio.sleep(0.02)

            # Emit completion
            yield f"data: {json.dumps({'type': 'execution_complete', 'data': {'execution_id': execution_id}})}\n\n"

        else:
            # Fallback: use reasoning orchestrator in agent mode
            yield f"data: {json.dumps({'type': 'step_update', 'data': {'step_id': 'reasoning', 'phase': 'planning', 'title': 'Analyzing your request...', 'status': 'active'}})}\n\n"

            full_content = ""
            async for event in reasoning_orchestrator.process_with_reasoning(
                query=query,
                session_id=request.session_id or str(uuid.uuid4()),
                mode="deep"
            ):
                event_type = event.get("type")
                event_data = event.get("data", {})

                if event_type == "layer_start":
                    yield f"data: {json.dumps({'type': 'step_update', 'data': {'step_id': event_data.get('layer_id', 'layer'), 'phase': event_data.get('type', 'research'), 'title': event_data.get('title', 'Processing...'), 'status': 'active'}})}\n\n"
                elif event_type == "layer_complete":
                    yield f"data: {json.dumps({'type': 'step_update', 'data': {'step_id': event_data.get('layer_id', 'layer'), 'phase': 'research', 'title': event_data.get('title', 'Complete'), 'status': 'complete'}})}\n\n"
                elif event_type == "sub_step":
                    yield f"data: {json.dumps({'type': 'execution_reasoning', 'data': {'chunk': event_data.get('step', '') + '\\n'}})}\n\n"
                elif event_type == "content":
                    chunk = event_data.get("chunk", "")
                    full_content += chunk
                    yield f"data: {json.dumps({'type': 'content', 'data': {'chunk': chunk}})}\n\n"
                elif event_type == "complete":
                    final = event_data.get("content", "")
                    if final and len(final) > len(full_content):
                        yield f"data: {json.dumps({'type': 'content', 'data': {'chunk': final[len(full_content):]  }})}\n\n"
                        full_content = final

                await asyncio.sleep(0.02)

            # Auto-detect tables for download
            table_lines = [l for l in full_content.split('\n') if '|' in l and not l.strip().startswith('---')]
            if len(table_lines) >= 3:
                try:
                    table_content = '\n'.join(table_lines)
                    file_path, filename, file_id = _generate_excel_from_content(table_content, "McLeuker AI Data")
                    yield f"data: {json.dumps({'type': 'download', 'data': {'filename': filename, 'download_url': f'/api/files/{file_id}.xlsx', 'file_id': f'{file_id}.xlsx', 'file_type': 'excel'}})}\n\n"
                except Exception as e:
                    logger.warning(f"Auto Excel generation failed: {e}")

            yield f"data: {json.dumps({'type': 'execution_complete', 'data': {'execution_id': execution_id}})}\n\n"

    except Exception as e:
        logger.error(f"Execution stream error: {e}", exc_info=True)
        yield f"data: {json.dumps({'type': 'execution_error', 'data': {'error': str(e)}})}\n\n"


@app.post("/api/v2/execute/stream")
async def execute_stream_v2(request: ExecuteRequest):
    """V2 SSE streaming execution endpoint — used by agent mode."""
    return StreamingResponse(
        _execute_stream_generator(request),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        }
    )


# Legacy alias
@app.post("/api/execute/stream")
async def execute_stream_legacy(request: ExecuteRequest):
    """Legacy alias for /api/v2/execute/stream."""
    return StreamingResponse(
        _execute_stream_generator(request),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        }
    )


@app.post("/api/execute")
async def execute_task(request: ExecuteRequest):
    """Non-streaming execution endpoint."""
    if not execution_engine:
        return JSONResponse(
            status_code=503,
            content={"success": False, "error": "Execution engine not available"}
        )

    try:
        results = []
        async for event in execution_engine.execute(
            query=request.get_query(),
            user_id=request.user_id,
            session_id=request.session_id,
            context=request.context,
        ):
            results.append(event)

        final_result = None
        for event in reversed(results):
            if event.get("type") == "result":
                final_result = event.get("data", {})
                break

        if final_result:
            return {"success": True, "result": final_result, "events": results}
        else:
            for event in reversed(results):
                if event.get("type") == "error":
                    return JSONResponse(
                        status_code=500,
                        content={"success": False, "error": event.get("data", {}).get("error", "Execution failed")}
                    )
            return {"success": True, "result": {"content": "Task completed", "events": results}}
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"success": False, "error": str(e)}
        )


@app.get("/api/v2/execute/{task_id}/status")
@app.get("/api/execute/{task_id}")
async def get_execution_status(task_id: str):
    """Get execution task status."""
    if not execution_engine:
        return JSONResponse(
            status_code=503,
            content={"success": False, "error": "Execution engine not available"}
        )
    task = execution_engine.get_task(task_id)
    if not task:
        return JSONResponse(
            status_code=404,
            content={"success": False, "error": "Task not found"}
        )
    return {"success": True, "task": task.to_dict()}


@app.get("/api/execute/active/list")
async def list_active_executions():
    if not execution_engine:
        return {"success": True, "tasks": []}
    return {"success": True, "tasks": execution_engine.get_active_tasks()}


@app.post("/api/v2/execute/{task_id}/pause")
async def pause_execution(task_id: str):
    return {"success": True, "message": f"Execution {task_id} paused"}


@app.post("/api/v2/execute/{task_id}/resume")
async def resume_execution(task_id: str):
    return {"success": True, "message": f"Execution {task_id} resumed"}


@app.post("/api/v2/execute/{task_id}/cancel")
async def cancel_execution(task_id: str):
    return {"success": True, "message": f"Execution {task_id} cancelled"}


@app.post("/api/v2/execute/credential")
async def submit_credential(data: dict):
    """Receive credentials from the user for execution."""
    return {"success": True, "message": "Credential received"}


# =============================================================================
# Nano Banana Image Generation
# =============================================================================

@app.post("/api/image/generate")
async def generate_image(request: ImageGenerateRequest):
    result = await nano_banana_service.generate_image(
        prompt=request.prompt,
        style=request.style,
        aspect_ratio=request.aspect_ratio
    )
    if result["success"]:
        return {
            "success": True,
            "image_data": result["image_data"],
            "mime_type": result["mime_type"],
            "prompt": result["prompt"]
        }
    return JSONResponse(
        status_code=400,
        content={"success": False, "error": result.get("error", "Image generation failed")}
    )


@app.post("/api/image/edit")
async def edit_image(request: ImageEditRequest):
    result = await nano_banana_service.edit_image(
        image_data=request.image_data,
        edit_prompt=request.edit_prompt,
        mime_type=request.mime_type
    )
    if result["success"]:
        return {
            "success": True,
            "image_data": result["image_data"],
            "mime_type": result["mime_type"],
            "edit_prompt": result["edit_prompt"]
        }
    return JSONResponse(
        status_code=400,
        content={"success": False, "error": result.get("error", "Image editing failed")}
    )


@app.post("/api/image/analyze")
async def analyze_image(request: ImageAnalyzeRequest):
    result = await nano_banana_service.analyze_image(
        image_data=request.image_data,
        question=request.question,
        mime_type=request.mime_type
    )
    if result["success"]:
        return {
            "success": True,
            "analysis": result["analysis"],
            "question": result["question"]
        }
    return JSONResponse(
        status_code=400,
        content={"success": False, "error": result.get("error", "Image analysis failed")}
    )


# =============================================================================
# File Upload & Download
# =============================================================================

@app.post("/api/upload")
async def upload_file(
    file: UploadFile = File(...),
    analyze: bool = Form(default=False)
):
    """Upload a file for processing."""
    try:
        content = await file.read()
        file_size = len(content)

        if file_size > 10 * 1024 * 1024:
            return JSONResponse(
                status_code=400,
                content={"success": False, "error": "File too large. Maximum size is 10MB."}
            )

        file_id = str(uuid.uuid4())[:8]
        content_type = file.content_type or "application/octet-stream"
        is_image = content_type.startswith("image/")
        base64_content = base64.b64encode(content).decode("utf-8")

        response = {
            "success": True,
            "file_id": file_id,
            "filename": file.filename,
            "content_type": content_type,
            "size": file_size,
            "base64_data": base64_content
        }

        if analyze and is_image:
            analysis_result = await nano_banana_service.analyze_image(
                image_data=base64_content,
                question="Describe this image in detail.",
                mime_type=content_type
            )
            if analysis_result["success"]:
                response["analysis"] = analysis_result["analysis"]

        return response

    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"success": False, "error": str(e)}
        )


@app.get("/api/files/{file_id}")
async def get_file(file_id: str):
    """Download generated files."""
    file_path = f"/tmp/mcleuker_files/{file_id}"

    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")

    media_types = {
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".csv": "text/csv",
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".md": "text/markdown",
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
    }

    ext = os.path.splitext(file_id)[1].lower()
    media_type = media_types.get(ext, "application/octet-stream")

    return FileResponse(
        file_path,
        media_type=media_type,
        filename=file_id,
        headers={"Content-Disposition": f'attachment; filename="{file_id}"'}
    )


# =============================================================================
# Document Generation (Legacy)
# =============================================================================

@app.post("/api/document/generate")
async def generate_document(request: DocumentGenerateRequest):
    """Generate a document from content (legacy endpoint)."""
    try:
        fmt = request.format.lower()
        if fmt in ("excel", "xlsx"):
            file_path, filename, file_id = _generate_excel_from_content(request.content, request.title)
            ext = "xlsx"
        elif fmt == "csv":
            file_path, filename, file_id = _generate_csv_from_content(request.content, request.title)
            ext = "csv"
        elif fmt == "pdf":
            file_path, filename, file_id = _generate_pdf_from_content(request.content, request.title)
            ext = "pdf"
        elif fmt in ("word", "docx"):
            file_path, filename, file_id = _generate_word_from_content(request.content, request.title)
            ext = "docx"
        elif fmt == "markdown":
            file_path, filename, file_id = _generate_markdown_from_content(request.content, request.title)
            ext = "md"
        elif fmt in ("powerpoint", "pptx"):
            file_path, filename, file_id = _generate_pptx_from_content(request.content, request.title)
            ext = "pptx"
        else:
            return JSONResponse(
                status_code=400,
                content={"success": False, "error": f"Unsupported format: {fmt}"}
            )

        return {
            "success": True,
            "file_id": f"{file_id}.{ext}",
            "filename": filename,
            "download_url": f"/api/files/{file_id}.{ext}",
            "format": fmt,
        }

    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"success": False, "error": str(e)}
        )


@app.post("/api/export/chat")
async def export_chat(request: ExportChatRequest):
    """Export chat conversation to a document."""
    try:
        content_lines = []
        for msg in request.messages:
            role = msg.get("role", "user")
            content = msg.get("content", "")
            if role == "user":
                content_lines.append(f"**User:** {content}")
            else:
                content_lines.append(f"**McLeuker AI:** {content}")
            content_lines.append("")

        content = "\n".join(content_lines)
        doc_request = DocumentGenerateRequest(
            content=content,
            title=request.title,
            format=request.format
        )
        return await generate_document(doc_request)

    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"success": False, "error": str(e)}
        )


# =============================================================================
# Capabilities & Status
# =============================================================================

@app.get("/api/capabilities")
async def get_capabilities():
    """Get the full capabilities of the McLeuker AI system."""
    return {
        "success": True,
        "capabilities": {
            "version": "10.0.0",
            "reasoning": {
                "multi_layer": True,
                "streaming": True,
                "modes": ["quick", "deep", "agent"],
            },
            "execution": {
                "available": execution_engine is not None,
                "web_automation": True,
                "file_generation": True,
                "code_execution": True,
                "api_calls": True,
                "credential_management": True,
                "supported_services": [
                    "github", "google", "twitter", "linkedin", "canva",
                    "slack", "notion", "figma", "vercel", "railway",
                ],
            },
            "image": {
                "generation": True,
                "editing": True,
                "analysis": True,
                "provider": "nano_banana",
            },
            "documents": {
                "formats": ["excel", "csv", "pdf", "word", "markdown", "pptx"],
                "auto_detect_tables": True,
                "export": True,
            },
            "domain_agents": {
                "available": enhancement is not None,
                "agents": [
                    "fashion", "beauty", "skincare", "sustainability",
                    "tech", "catwalk", "culture", "textile", "lifestyle",
                ] if enhancement else [],
            },
            "file_analysis": {
                "available": enhancement is not None,
                "supported_types": [
                    "image/jpeg", "image/png", "image/webp",
                    "application/pdf", "text/plain", "text/markdown",
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                ],
            },
        }
    }


@app.get("/api/enhancement/status")
async def get_enhancement_status():
    """Get the status of all enhancement components."""
    if not enhancement:
        return {
            "success": True,
            "status": {
                "enhancement_available": False,
                "execution_engine": execution_engine is not None,
            }
        }
    return {
        "success": True,
        "status": {
            "enhancement_available": True,
            "execution_engine": execution_engine is not None,
            **enhancement.get_status(),
        }
    }
