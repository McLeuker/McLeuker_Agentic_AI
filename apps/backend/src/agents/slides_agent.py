"""
Slides Agent — Professional Presentation Creation (PPTX)
=========================================================

Creates professional presentations from natural language descriptions.
Supports rich templates, layouts, charts, and visual elements.
"""

import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, AsyncGenerator, Dict, List, Optional
import openai

# Presentation library
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available - PPTX generation disabled")

logger = logging.getLogger(__name__)


@dataclass
class SlideContent:
    """Content for a single slide."""
    slide_type: str  # "title", "content", "two_column", "image", "chart", "section_header"
    title: str
    content: str = ""
    bullet_points: List[str] = field(default_factory=list)
    left_column: str = ""
    right_column: str = ""
    image_description: str = ""
    chart_data: Optional[Dict] = None
    speaker_notes: str = ""
    
    def to_dict(self) -> Dict:
        return {
            "slide_type": self.slide_type,
            "title": self.title,
            "content": self.content,
            "bullet_points": self.bullet_points,
            "left_column": self.left_column,
            "right_column": self.right_column,
            "image_description": self.image_description,
            "chart_data": self.chart_data,
            "speaker_notes": self.speaker_notes,
        }


@dataclass
class PresentationTemplate:
    """Template for presentation generation."""
    name: str
    description: str
    slide_structure: List[str]
    color_scheme: Dict[str, str]
    font_scheme: Dict[str, str]
    
    def to_dict(self) -> Dict:
        return {
            "name": self.name,
            "description": self.description,
            "slide_structure": self.slide_structure,
            "color_scheme": self.color_scheme,
            "font_scheme": self.font_scheme,
        }


@dataclass
class GeneratedPresentation:
    """A generated presentation."""
    presentation_id: str
    title: str
    file_path: str
    slides: List[SlideContent]
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    def to_dict(self) -> Dict:
        return {
            "presentation_id": self.presentation_id,
            "title": self.title,
            "file_path": self.file_path,
            "slides": [s.to_dict() for s in self.slides],
            "metadata": self.metadata,
        }


class SlidesAgent:
    """
    AI Slides Generator Agent for professional presentation creation.
    
    Usage:
        agent = SlidesAgent(llm_client)
        async for event in agent.generate_presentation(
            "Create a pitch deck for a SaaS startup"
        ):
            print(event)
    """
    
    # Built-in templates
    TEMPLATES = {
        "pitch_deck": PresentationTemplate(
            name="Pitch Deck",
            description="Startup pitch deck for investors",
            slide_structure=[
                "title",
                "problem",
                "solution",
                "market",
                "product",
                "traction",
                "business_model",
                "competition",
                "team",
                "financials",
                "ask",
                "closing",
            ],
            color_scheme={
                "primary": "#1a365d",
                "secondary": "#2b6cb0",
                "accent": "#48bb78",
                "background": "#ffffff",
                "text": "#1a202c",
            },
            font_scheme={
                "title": "Arial Bold",
                "body": "Arial",
            },
        ),
        "business_presentation": PresentationTemplate(
            name="Business Presentation",
            description="Professional business presentation",
            slide_structure=[
                "title",
                "agenda",
                "overview",
                "content_1",
                "content_2",
                "content_3",
                "data_analysis",
                "recommendations",
                "next_steps",
                "qa",
            ],
            color_scheme={
                "primary": "#2c5282",
                "secondary": "#4a5568",
                "accent": "#ed8936",
                "background": "#ffffff",
                "text": "#1a202c",
            },
            font_scheme={
                "title": "Calibri Bold",
                "body": "Calibri",
            },
        ),
        "training": PresentationTemplate(
            name="Training Presentation",
            description="Educational training presentation",
            slide_structure=[
                "title",
                "learning_objectives",
                "introduction",
                "concept_1",
                "concept_2",
                "concept_3",
                "examples",
                "exercise",
                "summary",
                "resources",
            ],
            color_scheme={
                "primary": "#276749",
                "secondary": "#38a169",
                "accent": "#d69e2e",
                "background": "#f7fafc",
                "text": "#1a202c",
            },
            font_scheme={
                "title": "Segoe UI Bold",
                "body": "Segoe UI",
            },
        ),
        "sales_deck": PresentationTemplate(
            name="Sales Deck",
            description="Sales presentation for prospects",
            slide_structure=[
                "title",
                "hook",
                "pain_points",
                "solution_overview",
                "features",
                "benefits",
                "case_study",
                "pricing",
                "implementation",
                "next_steps",
            ],
            color_scheme={
                "primary": "#c53030",
                "secondary": "#e53e3e",
                "accent": "#38b2ac",
                "background": "#ffffff",
                "text": "#1a202c",
            },
            font_scheme={
                "title": "Helvetica Bold",
                "body": "Helvetica",
            },
        ),
    }
    
    def __init__(
        self,
        llm_client: openai.AsyncOpenAI,
        model: str = "kimi-k2.5",
        output_dir: str = "/tmp/slides_agent",
    ):
        self.llm_client = llm_client
        self.model = model
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Track generated presentations
        self._presentations: Dict[str, GeneratedPresentation] = {}
    
    async def generate_presentation(
        self,
        description: str,
        template: Optional[str] = None,
        num_slides: Optional[int] = None,
        title: Optional[str] = None,
        context: Optional[Dict] = None,
    ) -> AsyncGenerator[Dict[str, Any], None]:
        """
        Generate a professional presentation from a description.
        
        Args:
            description: Natural language description of the presentation
            template: Optional template name to use
            num_slides: Target number of slides
            title: Optional presentation title
            context: Additional context for generation
            
        Yields:
            Progress events and final presentation info
        """
        import uuid
        presentation_id = str(uuid.uuid4())
        
        yield {"type": "start", "data": {"presentation_id": presentation_id, "description": description}}
        
        try:
            # Step 1: Analyze and select template
            yield {"type": "phase", "data": {"phase": "analysis", "status": "started"}}
            
            selected_template = None
            if template and template in self.TEMPLATES:
                selected_template = self.TEMPLATES[template]
            else:
                selected_template = await self._select_template(description)
            
            target_slides = num_slides or 10
            
            yield {
                "type": "phase",
                "data": {
                    "phase": "analysis",
                    "status": "completed",
                    "template": selected_template.name if selected_template else "custom",
                    "target_slides": target_slides,
                }
            }
            
            # Step 2: Generate presentation outline
            yield {"type": "phase", "data": {"phase": "outline", "status": "started"}}
            
            outline = await self._generate_outline(
                description,
                selected_template,
                target_slides,
                context,
            )
            
            yield {
                "type": "phase",
                "data": {
                    "phase": "outline",
                    "status": "completed",
                    "slides": len(outline),
                }
            }
            
            # Step 3: Generate slide content
            yield {"type": "phase", "data": {"phase": "content", "status": "started"}}
            
            slides = []
            for i, slide_info in enumerate(outline):
                slide = await self._generate_slide_content(
                    slide_info,
                    description,
                    selected_template,
                    context,
                )
                slides.append(slide)
                yield {"type": "progress", "data": {"slide": i + 1, "total": len(outline)}}
            
            yield {"type": "phase", "data": {"phase": "content", "status": "completed"}}
            
            # Step 4: Create PPTX file
            yield {"type": "phase", "data": {"phase": "creation", "status": "started"}}
            
            pres_title = title or await self._generate_title(description)
            file_path = await self._create_pptx(
                presentation_id,
                pres_title,
                slides,
                selected_template,
            )
            
            yield {"type": "phase", "data": {"phase": "creation", "status": "completed", "path": file_path}}
            
            # Step 5: Create presentation record
            presentation = GeneratedPresentation(
                presentation_id=presentation_id,
                title=pres_title,
                file_path=file_path,
                slides=slides,
                metadata={
                    "template": selected_template.name if selected_template else "custom",
                    "description": description,
                    "context": context or {},
                },
            )
            self._presentations[presentation_id] = presentation
            
            yield {
                "type": "complete",
                "data": {
                    "presentation_id": presentation_id,
                    "title": pres_title,
                    "file_path": file_path,
                    "slides": len(slides),
                }
            }
            
        except Exception as e:
            logger.error(f"Error generating presentation: {e}")
            yield {"type": "error", "data": {"message": str(e)}}
    
    async def _select_template(self, description: str) -> Optional[PresentationTemplate]:
        """Select the best template for the description."""
        messages = [
            {
                "role": "system",
                "content": f"""Select the best presentation template for this request.

Available templates:
{json.dumps({k: v.description for k, v in self.TEMPLATES.items()}, indent=2)}

Respond with JSON:
{{
    "template": "template_name or null for custom",
    "reasoning": "explanation"
}}"""
            },
            {
                "role": "user",
                "content": f"Select template for: {description}"
            }
        ]
        
        response = await self.llm_client.chat.completions.create(
            model=self.model,
            messages=messages,
            temperature=0.3,
            max_tokens=512,
            response_format={"type": "json_object"},
        )
        
        result = json.loads(response.choices[0].message.content)
        template_name = result.get("template")
        
        if template_name and template_name in self.TEMPLATES:
            return self.TEMPLATES[template_name]
        
        return None
    
    async def _generate_title(self, description: str) -> str:
        """Generate a title for the presentation."""
        messages = [
            {
                "role": "system",
                "content": "Generate a concise, engaging title for this presentation. Respond with just the title."
            },
            {
                "role": "user",
                "content": description
            }
        ]
        
        response = await self.llm_client.chat.completions.create(
            model=self.model,
            messages=messages,
            temperature=0.7,
            max_tokens=100,
        )
        
        return response.choices[0].message.content.strip().strip('"')
    
    async def _generate_outline(
        self,
        description: str,
        template: Optional[PresentationTemplate],
        num_slides: int,
        context: Optional[Dict],
    ) -> List[Dict]:
        """Generate presentation outline."""
        template_structure = template.slide_structure if template else []
        
        messages = [
            {
                "role": "system",
                "content": f"""Create a presentation outline with {num_slides} slides.

Template structure (if applicable): {json.dumps(template_structure)}

For each slide, specify:
- slide_type: "title", "content", "two_column", "image", "chart", "section_header"
- title: Slide title
- brief description of content

Respond with JSON:
{{
    "slides": [
        {{"slide_type": "type", "title": "Slide Title", "description": "Brief content description"}}
    ]
}}"""
            },
            {
                "role": "user",
                "content": f"Create outline for: {description}"
            }
        ]
        
        if context:
            messages.append({
                "role": "user",
                "content": f"Additional context: {json.dumps(context)}"
            })
        
        response = await self.llm_client.chat.completions.create(
            model=self.model,
            messages=messages,
            temperature=0.7,
            max_tokens=2048,
            response_format={"type": "json_object"},
        )
        
        result = json.loads(response.choices[0].message.content)
        return result.get("slides", [])
    
    async def _generate_slide_content(
        self,
        slide_info: Dict,
        description: str,
        template: Optional[PresentationTemplate],
        context: Optional[Dict],
    ) -> SlideContent:
        """Generate content for a single slide."""
        slide_type = slide_info.get("slide_type", "content")
        slide_title = slide_info.get("title", "")
        
        messages = [
            {
                "role": "system",
                "content": f"""Generate content for a presentation slide.

Slide type: {slide_type}
Slide title: {slide_title}
Presentation: {description}

Respond with JSON:
{{
    "title": "Slide title (can be different from above)",
    "content": "Main content paragraph (for content slides)",
    "bullet_points": ["point 1", "point 2", "point 3", "point 4"],
    "left_column": "Content for left column (two_column slides)",
    "right_column": "Content for right column (two_column slides)",
    "speaker_notes": "Notes for the presenter"
}}"""
            }
        ]
        
        response = await self.llm_client.chat.completions.create(
            model=self.model,
            messages=messages,
            temperature=0.7,
            max_tokens=1536,
            response_format={"type": "json_object"},
        )
        
        result = json.loads(response.choices[0].message.content)
        
        return SlideContent(
            slide_type=slide_type,
            title=result.get("title", slide_title),
            content=result.get("content", ""),
            bullet_points=result.get("bullet_points", []),
            left_column=result.get("left_column", ""),
            right_column=result.get("right_column", ""),
            speaker_notes=result.get("speaker_notes", ""),
        )
    
    async def _create_pptx(
        self,
        presentation_id: str,
        title: str,
        slides: List[SlideContent],
        template: Optional[PresentationTemplate],
    ) -> str:
        """Create a PPTX file."""
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not available")
        
        file_path = self.output_dir / f"{presentation_id}.pptx"
        
        prs = Presentation()
        
        # Set slide dimensions (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Get color scheme
        colors = template.color_scheme if template else {
            "primary": "#1a365d",
            "secondary": "#2b6cb0",
            "background": "#ffffff",
            "text": "#1a202c",
        }
        
        for slide_content in slides:
            # Select layout based on slide type
            if slide_content.slide_type == "title":
                layout = prs.slide_layouts[0]  # Title slide
            elif slide_content.slide_type == "section_header":
                layout = prs.slide_layouts[5]  # Blank with title
            else:
                layout = prs.slide_layouts[1]  # Title and content
            
            slide = prs.slides.add_slide(layout)
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = slide_content.title
            
            # Add content based on slide type
            if slide_content.slide_type == "content" and len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.text = slide_content.content
                
                # Add bullet points
                for point in slide_content.bullet_points:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
            
            elif slide_content.slide_type == "two_column":
                # Create two-column layout
                left = Inches(0.5)
                top = Inches(2)
                width = Inches(6)
                height = Inches(5)
                
                # Left column
                left_box = slide.shapes.add_textbox(left, top, width, height)
                tf = left_box.text_frame
                tf.text = slide_content.left_column
                
                # Right column
                right_box = slide.shapes.add_textbox(left + width + Inches(0.3), top, width, height)
                tf = right_box.text_frame
                tf.text = slide_content.right_column
            
            # Add speaker notes
            if slide_content.speaker_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_content.speaker_notes
        
        prs.save(str(file_path))
        return str(file_path)
    
    def list_templates(self) -> List[Dict]:
        """List available presentation templates."""
        return [{"id": k, **v.to_dict()} for k, v in self.TEMPLATES.items()]
    
    def get_presentation(self, presentation_id: str) -> Optional[GeneratedPresentation]:
        """Get a generated presentation by ID."""
        return self._presentations.get(presentation_id)


# Singleton instance
_slides_agent: Optional[SlidesAgent] = None


def get_slides_agent(llm_client: openai.AsyncOpenAI = None) -> Optional[SlidesAgent]:
    """Get or create the Slides Agent singleton."""
    global _slides_agent
    if _slides_agent is None and llm_client:
        _slides_agent = SlidesAgent(llm_client)
    return _slides_agent
